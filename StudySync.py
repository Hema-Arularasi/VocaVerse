import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize
import random
from pptx import Presentation
import fitz  # PyMuPDF
from gtts import gTTS
import os

nltk.download('punkt')
nltk.download('stopwords')

def extract_text_from_pdf(pdf_path):
    text = ""
    with fitz.open(pdf_path) as pdf_doc:
        for page_num in range(pdf_doc.page_count):
            page = pdf_doc[page_num]
            text += page.get_text("text")
    return text

def extract_text_from_ppt(ppt_path):
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def create_audio_file(text, audio_output_path):
    if text.strip():
        tts = gTTS(text, lang='en')
        tts.save(audio_output_path)
        print(f"Audio file with extracted words created and saved to {audio_output_path}.")
    else:
        print("No text to create an audio file.")

def ask_questions_about_theme(theme, para):
    print(f"\nHere are some questions about {theme}:\n")

    # Quiz questions
    print("Now, let's generate some quiz questions:\n")
    num_questions = int(input("Enter the number of quiz questions to generate: "))
    quiz_questions = generate_quiz_questions(para, num_questions)
    
    for i, question in enumerate(quiz_questions, start=1):
        print(f"Quiz Question {i}: {question}")

def generate_quiz_questions(text, num_questions=5):
    sentences = sent_tokenize(text)
    questions = []

    for _ in range(num_questions):
        random_sentence = random.choice(sentences)
        words = preprocess_text(random_sentence)
        
        if words:
            random_word = random.choice(words)
            question = f"What is the meaning of '{random_word}'?"
            questions.append(question)

    return questions

def preprocess_text(text):
    stop_words = set(stopwords.words("english"))
    words = [word.lower() for word in nltk.word_tokenize(text) if word.isalnum() and word.lower() not in stop_words]
    return words

def extract_info_from_file(file_path):
    if file_path.lower().endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        return extract_text_from_ppt(file_path)
    elif file_path.lower().endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    else:
        return None

def main():
    print("Welcome to the Note-Taking and Learning Assistant!")

    while True:
        user_input = input("Enter 'pdf', 'ppt', 'txt', 'topic', or 'exit': ").lower()

        if user_input == 'exit':
            print("Exiting the Note-Taking and Learning Assistant. Goodbye!")
            break
        elif user_input in ['pdf', 'ppt', 'txt']:
            file_path = input(f"Enter the path to the {user_input.upper()} file: ").strip()
            text = extract_info_from_file(file_path)

            # Print the extracted text
            print("Extracted Text:", text)

            # Ask user if they want to create an audio file for extracted words
            create_audio = input("Do you want to generate an audio file from the extracted words? (yes/no): ").lower()
            if create_audio == 'yes':
                audio_output_path = input("Enter the path to save the audio file (including filename and extension): ")
                create_audio_file(text, audio_output_path)

            if user_input == 'topic':
                ask_option = input(f"Do you want to answer questions about {user_input} content? (yes/no): ").lower()
                if ask_option == 'yes':
                    ask_questions_about_theme(user_input, text)  # Use user_input as the theme and file_path for extracting questions

        elif user_input == 'topic':
            theme = input("Enter the topic you want to learn about: ")
            para = input(f"Write a paragraph about {theme}: ")
            ask_questions_about_theme(user_input, para)  # Use user_input as the theme and para for extracting questions

            # Ask user if they want to create an audio file for extracted words
            create_audio = input("Do you want to generate an audio file from the extracted words? (yes/no): ").lower()
            if create_audio == 'yes':
                audio_output_path = input("Enter the path to save the audio file (including filename and extension): ")
                create_audio_file(para, audio_output_path)

        else:
            print("Invalid input. Please enter 'pdf', 'ppt', 'txt', 'topic', or 'exit.")

if __name__ == "__main__":
    main()